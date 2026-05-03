"""Convert final Markdown deliverables into DOCX, PDF, and PPTX."""

from __future__ import annotations

from pathlib import Path

from docx import Document
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGB
from pptx.util import Emu, Inches as PptxInches, Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

ROOT = Path(__file__).parent
OUT = ROOT / "report_results"

# ── Colour constants ─────────────────────────────────────────────────────────
_NAVY = (0x2E, 0x40, 0x57)          # dark-navy header background
_ROW_ALT = (0xE8, 0xEE, 0xF4)       # light-blue-grey alternate row


# ---------------------------------------------------------------------------
# Shared table-parsing helpers
# ---------------------------------------------------------------------------

def _split_table(lines: list[str]) -> list[list[str]]:
    """Parse Markdown table lines into a list of row cell lists.

    Separator rows (containing only ``-`` and ``|``) are discarded.

    Args:
        lines: Raw Markdown lines that all start with ``|``.

    Returns:
        List of rows; each row is a list of stripped cell strings.
    """
    rows = []
    for line in lines:
        if set(line.replace("|", "").strip()) <= {"-", " "}:
            continue
        rows.append([cell.strip() for cell in line.strip().strip("|").split("|")])
    return rows


def _parse_slide_table(body_lines: list[str]) -> tuple[list[list[str]], list[str]]:
    """Separate ``|``-prefixed table lines from regular body lines.

    Args:
        body_lines: Lines from one slide's body section.

    Returns:
        ``(table_rows, text_lines)`` where *table_rows* is the parsed table
        (empty list if no table found) and *text_lines* are the remaining
        non-table lines.
    """
    table_md: list[str] = []
    text_lines: list[str] = []
    for line in body_lines:
        if line.startswith("|"):
            table_md.append(line)
        else:
            text_lines.append(line)
    table_rows = _split_table(table_md) if table_md else []
    return table_rows, text_lines


# ---------------------------------------------------------------------------
# PPTX table helper
# ---------------------------------------------------------------------------

def _add_slide_table(slide, rows: list[list[str]]) -> None:  # type: ignore[type-arg]
    """Add a styled python-pptx table to *slide*.

    The first row is rendered as a dark-navy header with bold white text.
    Subsequent rows alternate between white and light-blue-grey fills.

    Args:
        slide: A python-pptx ``Slide`` object.
        rows: Parsed table rows (header first).
    """
    if not rows:
        return

    n_rows = len(rows)
    n_cols = len(rows[0])

    left = PptxInches(0.4)
    top = PptxInches(1.55)
    width = PptxInches(12.5)
    height = PptxInches(5.5)

    tbl = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    col_width = width // n_cols
    for col_idx in range(n_cols):
        tbl.columns[col_idx].width = col_width

    for r_idx, row_data in enumerate(rows):
        is_header = r_idx == 0
        fill_rgb = _NAVY if is_header else (_ROW_ALT if r_idx % 2 == 0 else (0xFF, 0xFF, 0xFF))
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = cell_text
            # Style the run
            run = cell.text_frame.paragraphs[0].runs
            if run:
                run[0].font.size = PptxPt(10)
                run[0].font.bold = is_header
                run[0].font.color.rgb = PptxRGB(0xFF, 0xFF, 0xFF) if is_header else PptxRGB(0x1A, 0x1A, 0x2E)
            # Fill the cell background
            from pptx.oxml.ns import qn as pqn
            import lxml.etree as etree
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, pqn("a:solidFill"))
            srgbClr = etree.SubElement(solidFill, pqn("a:srgbClr"))
            r, g, b = fill_rgb
            srgbClr.set("val", f"{r:02X}{g:02X}{b:02X}")


# ---------------------------------------------------------------------------
# DOCX conversion
# ---------------------------------------------------------------------------

def markdown_to_docx(md_path: Path, docx_path: Path) -> None:
    """Convert a Markdown report file to a Microsoft Word document.

    Handles headings (``#``, ``##``, ``###``), Markdown tables (with bold
    header row), code blocks, and appends speedup/efficiency chart images
    if present.

    Args:
        md_path: Path to the source ``.md`` file.
        docx_path: Destination path for the ``.docx`` output.
    """
    doc = Document()
    doc.core_properties.title = "Distributed Agentic GraphRAG Final Report"
    lines = md_path.read_text(encoding="utf-8").splitlines()
    idx = 0
    in_code = False
    while idx < len(lines):
        line = lines[idx]
        if line.startswith("```"):
            in_code = not in_code
            idx += 1
            continue
        if in_code:
            doc.add_paragraph(line, style="Intense Quote")
        elif line.startswith("### "):
            doc.add_heading(line[4:].strip(), level=2)
        elif line.startswith("## "):
            doc.add_heading(line[3:].strip(), level=1)
        elif line.startswith("# "):
            doc.add_heading(line[2:].strip(), level=0)
        elif line.startswith("|"):
            table_lines = []
            while idx < len(lines) and lines[idx].startswith("|"):
                table_lines.append(lines[idx])
                idx += 1
            rows = _split_table(table_lines)
            if rows:
                table = doc.add_table(rows=1, cols=len(rows[0]))
                table.style = "Table Grid"
                # Bold header row
                hdr_cells = table.rows[0].cells
                for col, value in enumerate(rows[0]):
                    hdr_cells[col].text = value
                    for para in hdr_cells[col].paragraphs:
                        for run in para.runs:
                            run.bold = True
                for row in rows[1:]:
                    cells = table.add_row().cells
                    for col, value in enumerate(row):
                        cells[col].text = value
            continue
        elif line.strip():
            doc.add_paragraph(line.strip())
        idx += 1

    for image_name in ["speedup_chart.png", "efficiency_chart.png"]:
        image_path = OUT / image_name
        if image_path.exists():
            doc.add_heading(image_name.replace("_", " ").replace(".png", "").title(), level=1)
            doc.add_picture(str(image_path), width=Inches(5.8))
    doc.save(docx_path)


# ---------------------------------------------------------------------------
# PDF conversion
# ---------------------------------------------------------------------------

def _rl_escape(text: str) -> str:
    """Escape ReportLab XML special characters in *text*."""
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def markdown_to_pdf(md_path: Path, pdf_path: Path) -> None:
    """Convert a Markdown report file to PDF using ReportLab.

    Handles headings, body paragraphs, Markdown tables (navy header,
    alternating rows), code blocks, and appends chart images if present.

    Args:
        md_path: Path to the source ``.md`` file.
        pdf_path: Destination path for the ``.pdf`` output.
    """
    styles = getSampleStyleSheet()
    styles.add(
        ParagraphStyle(
            name="CodeBlock",
            parent=styles["BodyText"],
            fontName="Courier",
            fontSize=8,
            leading=10,
        )
    )
    navy_color = colors.Color(_NAVY[0] / 255, _NAVY[1] / 255, _NAVY[2] / 255)
    alt_color = colors.Color(_ROW_ALT[0] / 255, _ROW_ALT[1] / 255, _ROW_ALT[2] / 255)

    story = []
    lines = md_path.read_text(encoding="utf-8").splitlines()
    idx = 0
    in_code = False
    while idx < len(lines):
        line = lines[idx]
        if line.startswith("```"):
            in_code = not in_code
            idx += 1
            continue
        if in_code:
            story.append(Paragraph(line.replace(" ", "&nbsp;"), styles["CodeBlock"]))
        elif line.startswith("### "):
            story.append(Paragraph(_rl_escape(line[4:].strip()), styles["Heading2"]))
        elif line.startswith("## "):
            story.append(Paragraph(_rl_escape(line[3:].strip()), styles["Heading1"]))
        elif line.startswith("# "):
            story.append(Paragraph(_rl_escape(line[2:].strip()), styles["Title"]))
            story.append(Spacer(1, 0.15 * inch))
        elif line.startswith("|"):
            table_lines = []
            while idx < len(lines) and lines[idx].startswith("|"):
                table_lines.append(lines[idx])
                idx += 1
            rows = _split_table(table_lines)
            if rows:
                # Escape special chars in every cell
                esc_rows = [[_rl_escape(c) for c in row] for row in rows]
                tbl = Table(esc_rows, repeatRows=1)
                ts_cmds = [
                    ("BACKGROUND", (0, 0), (-1, 0), navy_color),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                    ("FONT", (0, 1), (-1, -1), "Helvetica", 7),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ]
                # Alternating row fills (even indices starting at row 2)
                for r in range(1, len(rows)):
                    if r % 2 == 0:
                        ts_cmds.append(("BACKGROUND", (0, r), (-1, r), alt_color))
                tbl.setStyle(TableStyle(ts_cmds))
                story.append(tbl)
                story.append(Spacer(1, 0.15 * inch))
            continue
        elif line.strip():
            story.append(Paragraph(_rl_escape(line.strip()), styles["BodyText"]))
            story.append(Spacer(1, 0.08 * inch))
        idx += 1

    for image_name in ["speedup_chart.png", "efficiency_chart.png"]:
        image_path = OUT / image_name
        if image_path.exists():
            story.append(
                Paragraph(
                    image_name.replace("_", " ").replace(".png", "").title(),
                    styles["Heading1"],
                )
            )
            story.append(Image(str(image_path), width=5.6 * inch, height=3.5 * inch))
    SimpleDocTemplate(
        str(pdf_path),
        pagesize=letter,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36,
    ).build(story)


# ---------------------------------------------------------------------------
# PPTX conversion
# ---------------------------------------------------------------------------

def markdown_to_pptx(md_path: Path, pptx_path: Path) -> None:
    """Convert a slide Markdown file (``# Slide N: Title`` format) to PowerPoint.

    Each ``# Slide N:`` section becomes one slide.  Slides whose body contains
    ``|``-prefixed lines are rendered as native PPTX tables (navy header,
    alternating fills).  Chart images are embedded on the Results and PDC
    Analysis slides when the PNG files are present.

    Args:
        md_path: Path to the source slide ``.md`` file.
        pptx_path: Destination path for the ``.pptx`` output.
    """
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    text = md_path.read_text(encoding="utf-8")
    chunks = [chunk.strip() for chunk in text.split("# Slide ") if chunk.strip()]

    for chunk in chunks:
        lines = [line.strip() for line in chunk.splitlines() if line.strip()]
        if not lines:
            continue
        title = lines[0].split(":", 1)[1].strip() if ":" in lines[0] else lines[0]
        body = lines[1:]
        if title.lower() == "title" and body:
            title = body[0]
            body = body[1:]

        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title

        # Separate table rows from plain text lines
        table_rows, text_lines = _parse_slide_table(body)

        if table_rows:
            # Render as a native PPTX table; push the content placeholder off-screen
            ph = slide.shapes.placeholders[1]
            ph.left = PptxInches(0)
            ph.top = PptxInches(7.3)  # below slide boundary — effectively hidden
            ph.width = PptxInches(1)
            ph.height = PptxInches(0.1)
            ph.text_frame.text = ""

            # Add any non-table lines as small caption above table (if any)
            if text_lines:
                txBox = slide.shapes.add_textbox(
                    PptxInches(0.4), PptxInches(1.2), PptxInches(12.5), PptxInches(0.4)
                )
                for i, tl in enumerate(text_lines):
                    para = txBox.text_frame.paragraphs[0] if i == 0 else txBox.text_frame.add_paragraph()
                    para.text = tl
                    para.font.size = PptxPt(14)

            _add_slide_table(slide, table_rows)
        else:
            tf = slide.shapes.placeholders[1].text_frame
            tf.clear()
            for idx, item in enumerate(text_lines):
                paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                paragraph.text = item
                paragraph.font.size = PptxPt(24 if len(item) < 120 else 18)

        # Embed chart images by slide title keyword matching
        title_lower = title.lower()
        if any(kw in title_lower for kw in ("result", "experiment", "speedup")):
            img = OUT / "speedup_chart.png"
            if img.exists():
                slide.shapes.add_picture(
                    str(img), PptxInches(7.1), PptxInches(1.6), width=PptxInches(5.6)
                )
        if any(kw in title_lower for kw in ("pdc", "analysis", "efficiency", "parallel")):
            img = OUT / "efficiency_chart.png"
            if img.exists():
                slide.shapes.add_picture(
                    str(img), PptxInches(7.1), PptxInches(1.6), width=PptxInches(5.6)
                )

    prs.save(pptx_path)


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

def main() -> None:
    """Convert ``final_report.md`` and ``presentation_slides.md`` to DOCX, PDF, PPTX."""
    report_md = OUT / "final_report.md"
    slides_md = OUT / "presentation_slides.md"
    markdown_to_docx(report_md, OUT / "final_report.docx")
    markdown_to_pdf(report_md, OUT / "final_report.pdf")
    markdown_to_pptx(slides_md, OUT / "presentation_slides.pptx")
    print(f"Created: {OUT / 'final_report.docx'}")
    print(f"Created: {OUT / 'final_report.pdf'}")
    print(f"Created: {OUT / 'presentation_slides.pptx'}")


if __name__ == "__main__":
    main()
